import os
import logging
import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, Request

from aiogram import Bot, Dispatcher, F
from aiogram.enums import ParseMode
from aiogram.types import Message, FSInputFile, Update
from aiogram.client.default import DefaultBotProperties
from aiogram.fsm.storage.memory import MemoryStorage

from docx import Document
from pptx import Presentation

# Загрузка переменных окружения
load_dotenv()
BOT_TOKEN = os.getenv("BOT_TOKEN")
WEBHOOK_URL = os.getenv("WEBHOOK_URL")
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

WEBHOOK_PATH = f"/webhook/{BOT_TOKEN}"
FULL_WEBHOOK_URL = f"{WEBHOOK_URL}{WEBHOOK_PATH}"

# Логгирование
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("bot")

# Инициализация бота и диспетчера
bot = Bot(token=BOT_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher(storage=MemoryStorage())

# FastAPI
app = FastAPI()

# Стартовая команда
@dp.message(F.text == "/start")
async def start(message: Message):
    await message.answer(
        "👋 Привет! Я AI-помощник студентов.\n\n"
        "📄 Реферат: `реферат: тема`\n"
        "📊 Презентация: `презентация: тема`\n"
        "🧮 Задача: просто напиши её!"
    )

# Обработка всех сообщений
@dp.message()
async def handle_message(message: Message):
    try:
        text = message.text.lower()

        if text.startswith("реферат:"):
            prompt = text.split("реферат:", 1)[1].strip()
            await generate_docx(message, prompt)

        elif text.startswith("презентация:"):
            prompt = text.split("презентация:", 1)[1].strip()
            await generate_pptx(message, prompt)

        else:
            await generate_answer(message, message.text)

    except Exception as e:
        logger.error(f"Ошибка обработки сообщения: {e}")
        await message.answer("❌ Ошибка при обработке запроса.")

# Генерация обычного ответа
async def generate_answer(message: Message, prompt: str):
    try:
        response = await ask_openrouter(prompt)
        await message.answer(response)
    except Exception as e:
        logger.error(f"Ошибка генерации ответа: {e}")
        await message.answer("❌ Не удалось сгенерировать ответ.")

# Генерация реферата (.docx)
async def generate_docx(message: Message, prompt: str):
    try:
        content = await ask_openrouter(f"Напиши подробный реферат на тему: {prompt}")
        doc = Document()
        doc.add_heading(f"Реферат: {prompt}", 0)
        doc.add_paragraph(content)

        filename = f"ref_{message.chat.id}.docx"
        doc.save(filename)

        if os.path.exists(filename):
            await message.answer("📎 Отправляю реферат в формате .docx...")
            await message.answer_document(FSInputFile(filename))
            os.remove(filename)
        else:
            await message.answer("❌ Не удалось найти файл.")

    except Exception as e:
        logger.error(f"Ошибка при генерации реферата: {e}")
        await message.answer("❌ Ошибка при генерации реферата.")

# Генерация презентации (.pptx)
async def generate_pptx(message: Message, prompt: str):
    try:
        content = await ask_openrouter(f"Сделай структуру презентации по теме: {prompt}. Формат: Слайд 1: Заголовок - Описание")

        prs = Presentation()
        for line in content.split("\n"):
            if ":" in line:
                parts = line.split(":", 1)
                title = parts[0].strip()
                body = parts[1].strip()
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                slide.placeholders[1].text = body

        filename = f"ppt_{message.chat.id}.pptx"
        prs.save(filename)

        if os.path.exists(filename):
            await message.answer("📎 Отправляю презентацию в формате .pptx...")
            await message.answer_document(FSInputFile(filename))
            os.remove(filename)
        else:
            await message.answer("❌ Не удалось найти файл.")

    except Exception as e:
        logger.error(f"Ошибка при генерации презентации: {e}")
        await message.answer("❌ Ошибка при генерации презентации.")

# Запрос к OpenRouter API
async def ask_openrouter(prompt: str) -> str:
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "openai/gpt-4o-mini",
        "messages": [{"role": "user", "content": prompt}]
    }

    async with httpx.AsyncClient() as client:
        response = await client.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=data)

    result = response.json()
    if "choices" in result:
        return result["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Ошибка OpenRouter: {result}")

# Установка webhook
@app.on_event("startup")
async def on_startup():
    await bot.set_webhook(FULL_WEBHOOK_URL)
    logger.info(f"✅ Webhook установлен на {FULL_WEBHOOK_URL}")

@app.post(WEBHOOK_PATH)
async def process_webhook(request: Request):
    data = await request.json()
    update = Update.model_validate(data)
    await dp.feed_update(bot, update)
    return {"ok": True}
